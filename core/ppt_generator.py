import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def create_ppt_from_texts(texts: list[str], output_path: str = "output.pptx", skip_empty: bool = True):
    """
    Creates a PowerPoint presentation from a list of strings, 
    with one slide per string.
    """
    prs = Presentation()
    
    # Use a blank slide layout
    blank_slide_layout = prs.slide_layouts[6] 
    
    slide_count = 0

    for text in texts:
        # Skip empty texts if skip_empty is True
        if skip_empty and not text.strip():
            continue
            
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add a text box covering most of the slide
        left = top = Pt(50)
        width = prs.slide_width - Pt(100)
        height = prs.slide_height - Pt(100)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Vertically center the text box
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.text = text if text.strip() else ""
        
        # Horizontally center the paragraph
        p.alignment = PP_ALIGN.CENTER
        
        # Apply requested font and size
        p.font.name = 'Pretendard Light'
        p.font.size = Pt(66)
        
        # Alternate colors: Even slides = Black (#000000), Odd slides = Blue (#0070C0)
        if slide_count % 2 == 0:
            p.font.color.rgb = RGBColor(0, 0, 0)
        else:
            p.font.color.rgb = RGBColor(0, 112, 192)
            
        slide_count += 1
        
    prs.save(output_path)
    return output_path
